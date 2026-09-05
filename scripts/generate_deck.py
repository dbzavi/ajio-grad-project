from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_visual_deck():
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] 
    
    ajio_teal = RGBColor(0, 135, 122)
    ajio_dark = RGBColor(44, 62, 80)
    ajio_accent = RGBColor(255, 63, 108)
    bg_light = RGBColor(245, 247, 250)

    def add_header(slide, title_text):
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
        banner.fill.solid()
        banner.fill.fore_color.rgb = ajio_dark
        banner.line.color.rgb = ajio_dark
        tf = banner.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255,255,255)
        p.alignment = PP_ALIGN.CENTER

    # --- Slide 1: Overview (Based on Uploaded Screenshot) ---
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Title
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(3), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Overview"
    p.font.bold = True; p.font.size = Pt(28); p.font.color.rgb = RGBColor(255,0,0) # Red like screenshot

    # Header line
    slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(0.6), Inches(6), Inches(0.05)).fill.solid()

    # --- LEFT COLUMN (Width: 2.8, Left: 0.5) ---
    # AJIO Description
    desc = slide1.shapes.add_textbox(Inches(0.5), Inches(1), Inches(3), Inches(1.5))
    desc.text_frame.word_wrap = True
    p = desc.text_frame.paragraphs[0]
    p.text = "AJIO is a leading Indian fashion e-commerce platform offering a curated selection of apparel, footwear, and accessories. The platform focuses on providing trendy selections with a seamless browsing and wishlist experience."
    p.font.size = Pt(12)

    # Metrics Box (Peach background with hard numbers)
    metrics_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(3), Inches(1.5))
    metrics_box.fill.solid(); metrics_box.fill.fore_color.rgb = RGBColor(255, 230, 213); metrics_box.line.color.rgb = RGBColor(255, 230, 213)
    tf_m = metrics_box.text_frame
    
    # Row 1 (Numbers)
    p_m1 = tf_m.paragraphs[0]
    p_m1.text = "   30 Mn          48%          $12Bn"
    p_m1.font.size = Pt(16); p_m1.font.bold = True; p_m1.font.color.rgb = ajio_dark
    
    # Row 2 (Labels)
    p_m2 = tf_m.add_paragraph()
    p_m2.text = "Active Users     Abandonment     Market Opp."
    p_m2.font.size = Pt(10); p_m2.font.bold = False; p_m2.font.color.rgb = ajio_dark
    
    # Row 3 (Spacer)
    tf_m.add_paragraph().text = ""
    
    # Row 4 (Numbers)
    p_m3 = tf_m.add_paragraph()
    p_m3.text = "   18-35          12 min          60%+"
    p_m3.font.size = Pt(16); p_m3.font.bold = True; p_m3.font.color.rgb = ajio_dark
    
    # Row 5 (Labels)
    p_m4 = tf_m.add_paragraph()
    p_m4.text = "Average Age    Daily Session    Wishlist Usage"
    p_m4.font.size = Pt(10); p_m4.font.bold = False; p_m4.font.color.rgb = ajio_dark

    # Value Proposition Box
    vp_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(3), Inches(1.8))
    vp_box.fill.solid(); vp_box.fill.fore_color.rgb = RGBColor(255, 230, 213); vp_box.line.color.rgb = RGBColor(255, 230, 213)
    tf_vp = vp_box.text_frame
    p_vp1 = tf_vp.paragraphs[0]; p_vp1.text = "Value Proposition\n"; p_vp1.font.bold = True; p_vp1.font.size = Pt(14); p_vp1.font.color.rgb = RGBColor(255,0,0)
    p_vp2 = tf_vp.add_paragraph(); p_vp2.text = "• Trendy Fashion Selection\n• Easy Reverse Logistics\n• Seamless Wishlist curation\n• AI-driven Fit Recommendations"
    p_vp2.font.size = Pt(12); p_vp2.font.color.rgb = ajio_dark

    # Competitors Box
    comp_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6), Inches(3), Inches(1))
    comp_box.fill.solid(); comp_box.fill.fore_color.rgb = RGBColor(240, 240, 240); comp_box.line.color.rgb = RGBColor(240, 240, 240)
    tf_comp = comp_box.text_frame
    tf_comp.text = "Competitors:\nMyntra, Flipkart Fashion, Nykaa Fashion"
    tf_comp.paragraphs[0].font.size = Pt(12); tf_comp.paragraphs[0].font.color.rgb = RGBColor(255,0,0)

    # --- MIDDLE COLUMN (Width: 3.5, Left: 3.8) ---
    # Top Mission Box (Dashed)
    mission_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1), Inches(3.8), Inches(1))
    mission_box.fill.solid(); mission_box.fill.fore_color.rgb = RGBColor(255, 255, 255); mission_box.line.dash_style = 4 # Dashed
    tf_miss = mission_box.text_frame; tf_miss.word_wrap = True
    p_miss = tf_miss.paragraphs[0]
    p_miss.text = "🎯 Increase wishlist-to-purchase conversion rates among high-intent shoppers without utilizing restricted monetary incentives."
    p_miss.font.bold = True; p_miss.font.size = Pt(12); p_miss.font.color.rgb = ajio_dark

    # Why focus on this? Box (Dashed)
    focus_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(2.2), Inches(3.8), Inches(2.4))
    focus_box.fill.solid(); focus_box.fill.fore_color.rgb = RGBColor(255, 255, 255); focus_box.line.dash_style = 4
    tf_f = focus_box.text_frame; tf_f.word_wrap = True
    p_f1 = tf_f.paragraphs[0]; p_f1.text = "Why focus on this?\n"; p_f1.font.bold = True; p_f1.font.size = Pt(14); p_f1.font.color.rgb = RGBColor(255,0,0)
    p_f2 = tf_f.add_paragraph(); p_f2.text = "Predictable Margin – Bypassing monetary discounts preserves unit economics.\n\nImproved Data & Trust – AI sizing models refine future recommendations and build brand trust.\n\nStronger Brand Appeal – Solving a core industry pain point (Sizing Uncertainty) positions AJIO as a tech-forward leader."
    p_f2.font.size = Pt(11); p_f2.font.color.rgb = ajio_dark

    # Why do we focus on it now? Box (Dashed)
    now_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(4.8), Inches(3.8), Inches(2.5))
    now_box.fill.solid(); now_box.fill.fore_color.rgb = RGBColor(255, 255, 255); now_box.line.dash_style = 4
    tf_n = now_box.text_frame; tf_n.word_wrap = True
    p_n1 = tf_n.paragraphs[0]; p_n1.text = "Why do we focus on it now?\n"; p_n1.font.bold = True; p_n1.font.size = Pt(14); p_n1.font.color.rgb = RGBColor(255,0,0)
    p_n2 = tf_n.add_paragraph(); p_n2.text = "Market Maturity\nShoppers expect personalized, AI-driven experiences. LLM technology (RAG) is now cheap and fast enough to deploy at scale.\n\nHigh Abandonment\nSizing paranoia is currently the #1 non-monetary friction point blocking the checkout funnel.\n\nTrapped Revenue\nWishlists hold massive trapped GMV just waiting to be unlocked."
    p_n2.font.size = Pt(10); p_n2.font.color.rgb = ajio_dark

    # --- RIGHT COLUMN (Actors - Pink background) ---
    actors_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1), Inches(2), Inches(6.3))
    actors_box.fill.solid(); actors_box.fill.fore_color.rgb = RGBColor(255, 204, 204); actors_box.line.color.rgb = RGBColor(255, 204, 204)
    tf_a = actors_box.text_frame; tf_a.word_wrap = True
    
    p_a1 = tf_a.paragraphs[0]; p_a1.text = "Actors\n\n"; p_a1.font.bold = True; p_a1.font.size = Pt(14); p_a1.font.color.rgb = RGBColor(255,0,0)
    
    p_a2 = tf_a.add_paragraph(); p_a2.text = "Shoppers\nHesitant wishlisters who need sizing validation to complete purchases.\n\n"; p_a2.font.bold = True; p_a2.font.size = Pt(11); p_a2.font.color.rgb = ajio_dark
    p_a3 = tf_a.add_paragraph(); p_a3.text = "Fashion Brands\nPartners whose unit economics suffer from high size-related return rates.\n\n"; p_a3.font.bold = True; p_a3.font.size = Pt(11); p_a3.font.color.rgb = ajio_dark
    p_a4 = tf_a.add_paragraph(); p_a4.text = "AJIO Platform\nIncurs heavy reverse logistics costs from returns and loses out on trapped wishlist GMV."; p_a4.font.bold = True; p_a4.font.size = Pt(11); p_a4.font.color.rgb = ajio_dark

    # --- Slide 2: Business Outcome & KPI Tree (Based on Screenshot) ---
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Title (Blue chevron/ribbon style from screenshot)
    title_bg = slide2.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0), Inches(0.2), Inches(5), Inches(0.6))
    title_bg.fill.solid(); title_bg.fill.fore_color.rgb = RGBColor(218, 227, 243); title_bg.line.color.rgb = ajio_dark
    tf_t2 = title_bg.text_frame
    p_t2 = tf_t2.paragraphs[0]; p_t2.text = "Understanding the Business Outcome to target"; p_t2.font.bold = True; p_t2.font.size = Pt(16); p_t2.font.color.rgb = ajio_dark

    # --- TOP LEFT: KPI Tree Diagram (Background Box) ---
    tree_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1), Inches(6.5), Inches(3.5))
    tree_bg.fill.solid(); tree_bg.fill.fore_color.rgb = RGBColor(248, 249, 250); tree_bg.line.color.rgb = RGBColor(248, 249, 250)

    # Root Node
    root_node = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(2.3), Inches(1.8), Inches(0.6))
    root_node.fill.solid(); root_node.fill.fore_color.rgb = RGBColor(255,255,255); root_node.line.color.rgb = RGBColor(200,200,200)
    tf_r = root_node.text_frame; tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "Wishlist-to-Purchase\nConversion Rate"
    tf_r.paragraphs[0].font.bold = True; tf_r.paragraphs[0].font.size = Pt(11); tf_r.paragraphs[0].font.color.rgb = ajio_dark; tf_r.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Branch Nodes
    b1 = slide2.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(2), Inches(0.5)); b1.text_frame.text = "# Users Visiting Wishlist"
    b2 = slide2.shapes.add_textbox(Inches(2.5), Inches(2.3), Inches(2), Inches(0.5)); b2.text_frame.text = "Checkout Initiation Rate"
    b3 = slide2.shapes.add_textbox(Inches(2.5), Inches(3.3), Inches(2), Inches(0.5)); b3.text_frame.text = "Checkout Completion Rate"
    for b in [b1, b2, b3]: b.text_frame.paragraphs[0].font.size = Pt(10); b.text_frame.paragraphs[0].font.bold = True

    # Sub-Branch Nodes
    sb1 = slide2.shapes.add_textbox(Inches(4.8), Inches(1.1), Inches(1.8), Inches(0.5)); sb1.text_frame.text = "App Opens\n✖ Wishlist CTR"
    sb2 = slide2.shapes.add_textbox(Inches(4.8), Inches(2.1), Inches(1.8), Inches(0.5)); sb2.text_frame.text = "Sizing Confidence\n✖ Price Acceptance"
    sb3 = slide2.shapes.add_textbox(Inches(4.8), Inches(3.1), Inches(1.8), Inches(0.5)); sb3.text_frame.text = "Payment Success\n✖ Address Drop-off"
    for sb in [sb1, sb2, sb3]: sb.text_frame.paragraphs[0].font.size = Pt(9); sb.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,100,100)

    # Drawing Connecting Lines (Simulated with thin rectangles)
    # Root to vertical spine
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(2.6), Inches(0.2), Inches(0.02)).fill.solid()
    # Vertical Spine
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(1.5), Inches(0.02), Inches(2.0)).fill.solid()
    # Horizontal connectors to branches
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(1.5), Inches(0.1), Inches(0.02)).fill.solid()
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(2.5), Inches(0.1), Inches(0.02)).fill.solid()
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(3.5), Inches(0.1), Inches(0.02)).fill.solid()
    
    # Plus symbols between branches
    p1 = slide2.shapes.add_textbox(Inches(2.3), Inches(1.8), Inches(0.5), Inches(0.5)); p1.text_frame.text = "➕"; p1.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    p2 = slide2.shapes.add_textbox(Inches(2.3), Inches(2.8), Inches(0.5), Inches(0.5)); p2.text_frame.text = "➕"; p2.text_frame.paragraphs[0].font.color.rgb = ajio_dark

    # --- BOTTOM LEFT: Outcomes & Insights ---
    # Leading/Lagging Box
    outcomes = slide2.shapes.add_textbox(Inches(0.2), Inches(4.8), Inches(3), Inches(2.5))
    tf_out = outcomes.text_frame
    p_o1 = tf_out.paragraphs[0]; p_o1.text = "Product Outcome (Leading) :"; p_o1.font.bold = True; p_o1.font.size = Pt(12)
    p_o2 = tf_out.add_paragraph(); p_o2.text = "• Fit-Match Widget Engagement Rate\n• Time-to-Purchase (Wishlist -> Cart)"; p_o2.font.size = Pt(11)
    tf_out.add_paragraph().text = ""
    p_o3 = tf_out.add_paragraph(); p_o3.text = "Business Outcome (Lagging) :"; p_o3.font.bold = True; p_o3.font.size = Pt(12)
    p_o4 = tf_out.add_paragraph(); p_o4.text = "• Wishlist-to-Purchase Conv. Rate\n• Size-Related Return Rate"; p_o4.font.size = Pt(11)

    # Key Insights Box
    insights = slide2.shapes.add_textbox(Inches(3.3), Inches(4.8), Inches(3.4), Inches(2.5))
    tf_ins = insights.text_frame; tf_ins.word_wrap = True
    p_i1 = tf_ins.paragraphs[0]; p_i1.text = "Key Insights"; p_i1.font.bold = True; p_i1.font.size = Pt(14); p_i1.font.color.rgb = ajio_dark
    p_i2 = tf_ins.add_paragraph(); p_i2.text = "• Users treat the wishlist as a holding area while they seek external validation for sizing.\n• Solving the 'visual and data trust gap' directly unlocks checkout initiation without requiring restricted monetary discounts."; p_i2.font.size = Pt(11)

    # Vertical divider line
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(4.8), Inches(0.01), Inches(2)).fill.solid()

    # --- RIGHT COLUMN ---
    # Target Box
    target_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(0.2), Inches(2.9), Inches(1))
    target_bg.fill.solid(); target_bg.fill.fore_color.rgb = RGBColor(218, 227, 243); target_bg.line.color.rgb = ajio_dark
    tf_targ = target_bg.text_frame
    p_tg1 = tf_targ.paragraphs[0]; p_tg1.text = "🎯 Target Business Outcome\n"; p_tg1.font.bold = True; p_tg1.font.size = Pt(12); p_tg1.font.color.rgb = ajio_dark; p_tg1.alignment = PP_ALIGN.CENTER
    p_tg2 = tf_targ.add_paragraph(); p_tg2.text = "Increase Wishlist-to-Purchase Rate"; p_tg2.font.bold = True; p_tg2.font.size = Pt(11); p_tg2.alignment = PP_ALIGN.CENTER

    # Paragraph Description
    desc_rt = slide2.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(3), Inches(3))
    tf_drt = desc_rt.text_frame; tf_drt.word_wrap = True
    p_drt = tf_drt.paragraphs[0]
    p_drt.text = "The overall wishlist conversion rate can be increased by providing extreme sizing confidence directly on the product page.\n\nSince price discounts are restricted by constraints, addressing sizing paranoia is the highest-leverage intervention to transition users from 'evaluation' to 'checkout'."
    p_drt.font.size = Pt(11)

    # Actors Box
    act_rt = slide2.shapes.add_textbox(Inches(6.8), Inches(4.2), Inches(3), Inches(3))
    tf_art = act_rt.text_frame; tf_art.word_wrap = True
    p_a_t = tf_art.paragraphs[0]; p_a_t.text = "Actors in the System"; p_a_t.font.bold = True; p_a_t.font.size = Pt(14); p_a_t.font.color.rgb = ajio_dark
    p_a_1 = tf_art.add_paragraph(); p_a_1.text = "1. Hesitant Wishlisters: High intent, low sizing confidence.\n2. AJIO Platform: Manages checkout flow.\n3. RAG AI: Synthesizes reviews."; p_a_1.font.size = Pt(11)

    # --- Slide 3: User Research & Persona (Based on Screenshot) ---
    slide3 = prs.slides.add_slide(blank_layout)
    
    # Title (Blue chevron style)
    title_bg3 = slide3.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0), Inches(0.1), Inches(6), Inches(0.5))
    title_bg3.fill.solid(); title_bg3.fill.fore_color.rgb = RGBColor(218, 227, 243); title_bg3.line.color.rgb = ajio_dark
    tf_t3 = title_bg3.text_frame
    p_t3 = tf_t3.paragraphs[0]; p_t3.text = "Understanding our Shoppers and their pain points"; p_t3.font.bold = True; p_t3.font.size = Pt(14); p_t3.font.color.rgb = ajio_dark

    # Top Right Evidence Box
    ev_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0.1), Inches(3.2), Inches(0.5))
    ev_box.fill.solid(); ev_box.fill.fore_color.rgb = RGBColor(255, 255, 255); ev_box.line.color.rgb = ajio_dark
    p_ev = ev_box.text_frame.paragraphs[0]; p_ev.text = "30 Surveys & 1,000+ AI Analyzed Reviews"; p_ev.font.size = Pt(10); p_ev.font.color.rgb = ajio_dark; p_ev.alignment = PP_ALIGN.CENTER

    # Footer Breadcrumbs (Bottom)
    footer = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(10), Inches(0.4))
    footer.fill.solid(); footer.fill.fore_color.rgb = ajio_dark; footer.line.color.rgb = ajio_dark
    tf_f = footer.text_frame
    p_f = tf_f.paragraphs[0]; p_f.text = "Product & Goal   |   USER RESEARCH   |   Problem Definition   |   Solution Ideation   |   Metrics"
    p_f.font.size = Pt(10); p_f.font.color.rgb = RGBColor(255,255,255); p_f.alignment = PP_ALIGN.CENTER

    # --- LEFT COLUMN (Width 4.5, Left 0.2) ---
    left_col = slide3.shapes.add_textbox(Inches(0.2), Inches(0.8), Inches(4.5), Inches(6))
    tf_lc = left_col.text_frame; tf_lc.word_wrap = True

    # Target User Segment
    p_lc1 = tf_lc.paragraphs[0]; p_lc1.text = "Target User Segment"; p_lc1.font.bold = True; p_lc1.font.size = Pt(12); p_lc1.font.color.rgb = ajio_dark
    p_lc2 = tf_lc.add_paragraph(); p_lc2.text = "Hesitant Wishlister (Gen-Z/Millennials, Ages 18-35). Highly engaged shoppers who curate wishlists but hesitate at checkout due to fit uncertainty."; p_lc2.font.size = Pt(10)
    tf_lc.add_paragraph().text = ""

    # Market Size
    p_lc3 = tf_lc.add_paragraph(); p_lc3.text = "Estimated Target Market Size"; p_lc3.font.bold = True; p_lc3.font.size = Pt(12); p_lc3.font.color.rgb = ajio_dark
    p_lc4 = tf_lc.add_paragraph(); p_lc4.text = "• Active AJIO Users ≈ 30M\n• High-intent wishlisters ≈ 15M (50% of active)\n• Cart abandonment due to sizing (48%) ≈ 7.2M users impacted directly."; p_lc4.font.size = Pt(10)
    tf_lc.add_paragraph().text = ""

    # Pain Points
    p_lc5 = tf_lc.add_paragraph(); p_lc5.text = "Consolidated Top Pain Points"; p_lc5.font.bold = True; p_lc5.font.size = Pt(12); p_lc5.font.color.rgb = ajio_dark
    p_lc6 = tf_lc.add_paragraph(); p_lc6.text = "• Sizing inconsistency across brands.\n• Lack of visual representation on average bodies.\n• Fear of tedious reverse logistics for returns."; p_lc6.font.size = Pt(10)
    
    # JTBD divider
    slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(4.5), Inches(4.5), Inches(0.01)).fill.solid()
    
    # JTBD
    jtbd = slide3.shapes.add_textbox(Inches(0.2), Inches(4.6), Inches(4.5), Inches(2))
    tf_j = jtbd.text_frame; tf_j.word_wrap = True
    p_j1 = tf_j.paragraphs[0]; p_j1.text = "JTBD - Jobs To Be Done"; p_j1.font.bold = True; p_j1.font.size = Pt(12); p_j1.font.color.rgb = ajio_dark
    p_j2 = tf_j.add_paragraph(); p_j2.text = "When I am shopping for new styles on AJIO,\nBut I struggle to confidently map the generic size chart to my unique body type,\nPlease help me visualize and validate the fit before purchasing,\nSo that I can buy with confidence and avoid the hassle of returning ill-fitting clothes."; p_j2.font.size = Pt(10)

    # --- RIGHT COLUMN (Width 4.8, Left 4.9) ---
    right_col = slide3.shapes.add_textbox(Inches(4.9), Inches(0.8), Inches(4.8), Inches(2.5))
    tf_rc = right_col.text_frame; tf_rc.word_wrap = True
    
    # Insights Title
    p_rc1 = tf_rc.paragraphs[0]; p_rc1.text = "Quantitative Survey Data (N=30)"; p_rc1.font.bold = True; p_rc1.font.size = Pt(14); p_rc1.font.color.rgb = ajio_dark
    
    # Chart 1: Reason for Delay
    lbl1 = slide3.shapes.add_textbox(Inches(5.0), Inches(1.3), Inches(4), Inches(0.3))
    lbl1.text_frame.paragraphs[0].text = "Q: What is the #1 reason you delay buying wishlisted items?"; lbl1.text_frame.paragraphs[0].font.size = Pt(10)
    
    # 76.7% Bar
    bar1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.6), Inches(3.4), Inches(0.3))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = RGBColor(239, 68, 68); bar1.line.color.rgb = RGBColor(239, 68, 68)
    bar1.text_frame.paragraphs[0].text = " 76.7% - Sizing/Fit Uncertainty"; bar1.text_frame.paragraphs[0].font.size = Pt(10); bar1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    
    # 23.3% Bar
    bar2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.95), Inches(1.0), Inches(0.3))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(200, 200, 200); bar2.line.color.rgb = RGBColor(200, 200, 200)
    bar2.text_frame.paragraphs[0].text = " 23.3% - Price"; bar2.text_frame.paragraphs[0].font.size = Pt(9); bar2.text_frame.paragraphs[0].font.color.rgb = ajio_dark

    # Chart 2: External Validation
    lbl2 = slide3.shapes.add_textbox(Inches(5.0), Inches(2.4), Inches(4), Inches(0.3))
    lbl2.text_frame.paragraphs[0].text = "Q: Do you leave the app to find sizing reviews on other platforms?"; lbl2.text_frame.paragraphs[0].font.size = Pt(10)

    # 90% Bar
    bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(2.7), Inches(4.0), Inches(0.3))
    bar3.fill.solid(); bar3.fill.fore_color.rgb = ajio_dark; bar3.line.color.rgb = ajio_dark
    bar3.text_frame.paragraphs[0].text = " 90.0% - Yes, actively hunt for reviews"; bar3.text_frame.paragraphs[0].font.size = Pt(10); bar3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    # Persona Box Background (Light blue)
    persona_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(3.5), Inches(4.5), Inches(3.4))
    persona_bg.fill.solid(); persona_bg.fill.fore_color.rgb = RGBColor(226, 235, 248); persona_bg.line.color.rgb = ajio_dark
    
    # Persona Tab (Vertical text simulation - rotating a text box)
    tab = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(4.2), Inches(0.8), Inches(2))
    tab.fill.solid(); tab.fill.fore_color.rgb = RGBColor(180, 200, 230); tab.line.color.rgb = ajio_dark
    tab.rotation = 270 # Vertical
    p_tab = tab.text_frame.paragraphs[0]; p_tab.text = "Persona - Wishlister"; p_tab.font.size = Pt(9); p_tab.alignment = PP_ALIGN.CENTER

    # Persona Content
    persona_txt = slide3.shapes.add_textbox(Inches(5.4), Inches(3.6), Inches(4.3), Inches(3.2))
    tf_p = persona_txt.text_frame; tf_p.word_wrap = True
    
    p_p1 = tf_p.paragraphs[0]
    p_p1.text = "Priya | 26 Years | Software Engineer | Bangalore"
    p_p1.font.bold = True; p_p1.font.size = Pt(12); p_p1.font.color.rgb = ajio_dark
    
    p_p2 = tf_p.add_paragraph()
    p_p2.text = "26-year-old professional looking to upgrade her wardrobe, but hesitant to buy online due to past experiences with poor-fitting garments."
    p_p2.font.size = Pt(10); p_p2.font.italic = True
    
    p_p3 = tf_p.add_paragraph(); p_p3.text = "\nNeeds"; p_p3.font.bold = True; p_p3.font.size = Pt(11); p_p3.font.color.rgb = ajio_dark
    p_p4 = tf_p.add_paragraph(); p_p4.text = "• Visual validation of fit (Virtual Try-On)\n• Community consensus on sizing (runs large/small)"; p_p4.font.size = Pt(10)
    
    p_p5 = tf_p.add_paragraph(); p_p5.text = "Pain Points"; p_p5.font.bold = True; p_p5.font.size = Pt(11); p_p5.font.color.rgb = RGBColor(200,50,50)
    p_p6 = tf_p.add_paragraph(); p_p6.text = "• Uses the wishlist purely as a bookmark while hunting for external social proof (YouTube, Reddit).\n• Abandons cart if external reviews aren't found."; p_p6.font.size = Pt(10)

    # --- Slide 4: Problem Statement (Problem Framing Canvas) ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Problem Framing Canvas")
    
    # 1. What is the true problem? (Top Left, Width: 3, Height: 3)
    b1 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1.2), Inches(3.1), Inches(3.3))
    b1.fill.solid(); b1.fill.fore_color.rgb = RGBColor(255,255,255); b1.line.color.rgb = ajio_dark
    tf1 = b1.text_frame; tf1.word_wrap = True
    p1 = tf1.paragraphs[0]; p1.text = "What is the true problem?\n"; p1.font.bold = True; p1.font.size = Pt(14); p1.font.color.rgb = ajio_dark
    p1_2 = tf1.add_paragraph(); p1_2.text = "The 'Sizing Trust Deficit': A behavioral gap where users add to wishlist but abandon the cart due to:\n"
    p1_3 = tf1.add_paragraph(); p1_3.text = "• Lack of visual representation on average bodies.\n• Fear of the tedious returns process.\n• Inability to mentally map generic sizes to their unique body type."
    for p in [p1_2, p1_3]: p.font.size = Pt(12); p.font.color.rgb = ajio_dark

    # 2. Who are the customers facing the problem? (Top Middle, Width: 3.1, Height: 3)
    b2 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.45), Inches(1.2), Inches(3.1), Inches(3.3))
    b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(255,255,255); b2.line.color.rgb = ajio_dark
    tf2 = b2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = "Who are the customers facing the problem?\n"; p2.font.bold = True; p2.font.size = Pt(14); p2.font.color.rgb = ajio_dark
    p2_2 = tf2.add_paragraph(); p2_2.text = "The 'Hesitant Wishlisters' (Ages 20-35) in the AJIO app.\n"
    p2_3 = tf2.add_paragraph(); p2_3.text = "• Frequent shoppers who curate high-intent wishlists but suffer from low checkout completion.\n• Users who rely heavily on social proof and external validation before committing."
    for p in [p2_2, p2_3]: p.font.size = Pt(12); p.font.color.rgb = ajio_dark

    # 3. How do we know it is a real problem? (Top Right, Width: 3.1, Height: 3)
    b3 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(1.2), Inches(3.1), Inches(3.3))
    b3.fill.solid(); b3.fill.fore_color.rgb = RGBColor(255,255,255); b3.line.color.rgb = ajio_dark
    tf3 = b3.text_frame; tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.text = "How do we know it is a real problem?\n"; p3.font.bold = True; p3.font.size = Pt(14); p3.font.color.rgb = ajio_dark
    p3_2 = tf3.add_paragraph(); p3_2.text = "• Evidence: Survey shows 76.7% state 'Uncertainty about size/fit' is the #1 reason they delay buying.\n"
    p3_3 = tf3.add_paragraph(); p3_3.text = "• Friction: 90% actively leave the app to hunt for reviews from buyers with similar bodies.\n"
    p3_4 = tf3.add_paragraph(); p3_4.text = "• AI Analysis: Sizing inconsistency drives 48% of non-monetary cart abandonment."
    for p in [p3_2, p3_3, p3_4]: p.font.size = Pt(12); p.font.color.rgb = ajio_dark

    # 4. What is the value generated? (Bottom Left, Width: 6, Height: 2.5)
    b4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(4.7), Inches(6.35), Inches(2.6))
    b4.fill.solid(); b4.fill.fore_color.rgb = RGBColor(255,255,255); b4.line.color.rgb = ajio_dark
    tf4 = b4.text_frame; tf4.word_wrap = True
    p4 = tf4.paragraphs[0]; p4.text = "What is the value generated by solving this problem?\n"; p4.font.bold = True; p4.font.size = Pt(14); p4.font.color.rgb = ajio_dark
    
    p4_2 = tf4.add_paragraph(); p4_2.text = "FOR THE TARGET CUSTOMERS:\n• Eliminating the 'External Review-Hunting' time tax.\n• Providing extreme purchasing confidence through personalized data.\n"; p4_2.font.size = Pt(11); p4_2.font.color.rgb = ajio_dark
    p4_3 = tf4.add_paragraph(); p4_3.text = "FOR THE BUSINESS:\n• Direct increase in the Wishlist-to-Purchase Conversion Rate (North Star).\n• Bypassing the need to offer restrictive monetary discounts.\n• Drastic reduction in expensive reverse-logistics costs (Size-Related Returns)."; p4_3.font.size = Pt(11); p4_3.font.color.rgb = ajio_dark

    # 5. Why should we solve this problem now? (Bottom Right, Width: 3.5, Height: 2.5)
    b5 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(4.7), Inches(3.1), Inches(2.6))
    b5.fill.solid(); b5.fill.fore_color.rgb = RGBColor(255,255,255); b5.line.color.rgb = ajio_dark
    tf5 = b5.text_frame; tf5.word_wrap = True
    p5 = tf5.paragraphs[0]; p5.text = "Why should we solve this problem now?\n"; p5.font.bold = True; p5.font.size = Pt(14); p5.font.color.rgb = ajio_dark
    p5_2 = tf5.add_paragraph(); p5_2.text = "• Market Maturity: E-commerce is shifting to personalized, AI-driven experiences.\n"
    p5_3 = tf5.add_paragraph(); p5_3.text = "• Growth: Wishlists hold massive trapped revenue; unlocking it drives immediate ROI.\n"
    p5_4 = tf5.add_paragraph(); p5_4.text = "• Platform War: Dominating 'Fit Trust' builds a moat against competitors relying purely on discounts."
    for p in [p5_2, p5_3, p5_4]: p.font.size = Pt(11); p.font.color.rgb = ajio_dark

    # --- Slide 5: Ideating Possible Solutions (Based on Screenshot) ---
    slide5 = prs.slides.add_slide(blank_layout)
    
    # Title
    t_bg5 = slide5.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(4), Inches(0.5))
    t_bg5.text_frame.paragraphs[0].text = "Ideating possible solutions"; t_bg5.text_frame.paragraphs[0].font.bold = True; t_bg5.text_frame.paragraphs[0].font.size = Pt(20); t_bg5.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Red Line extending from title
    slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(0.4), Inches(5.2), Inches(0.03)).fill.solid()

    # Column 1: AJIO Fit-Match AI (Our winning solution)
    col1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.8), Inches(3.1), Inches(6.3))
    col1.fill.solid(); col1.fill.fore_color.rgb = RGBColor(254, 240, 230); col1.line.color.rgb = ajio_dark; col1.line.dash_style = 4
    
    # Col 1 Header
    h1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.9), Inches(2.7), Inches(0.4))
    h1.fill.solid(); h1.fill.fore_color.rgb = RGBColor(239, 68, 68); h1.line.color.rgb = RGBColor(239, 68, 68)
    h1.text_frame.paragraphs[0].text = "AJIO Fit-Match AI"; h1.text_frame.paragraphs[0].font.bold = True; h1.text_frame.paragraphs[0].font.size = Pt(12); h1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; h1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    
    # Col 1 Text
    tf1 = col1.text_frame; tf1.word_wrap = True
    tf1.add_paragraph().text = "\n\n🔹 Problem: Users cannot map generic sizes to their unique bodies and hate hunting for reviews.\n\n🔹 Solution: A RAG-powered AI widget that synthesizes community reviews to provide extreme sizing confidence.\n\nKey Features:\n✅ Review Synthesis – Instantly summarizes 1,000+ reviews.\n✅ Fit Prediction – Recommends exact size based on user height/weight.\n✅ Virtual Try-On – Visually maps the garment.\n\n\n⭐ Impact: Unlocks immediate checkout by bridging both data and visual trust gaps without discounting."
    for p in tf1.paragraphs: p.font.size = Pt(11); p.font.color.rgb = ajio_dark

    # Column 2: AJIO Size-Swap Guarantee
    col2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.45), Inches(0.8), Inches(3.1), Inches(6.3))
    col2.fill.solid(); col2.fill.fore_color.rgb = RGBColor(254, 240, 230); col2.line.color.rgb = ajio_dark; col2.line.dash_style = 4
    
    # Col 2 Header
    h2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.65), Inches(0.9), Inches(2.7), Inches(0.4))
    h2.fill.solid(); h2.fill.fore_color.rgb = RGBColor(239, 68, 68); h2.line.color.rgb = RGBColor(239, 68, 68)
    h2.text_frame.paragraphs[0].text = "AJIO Size-Swap Guarantee"; h2.text_frame.paragraphs[0].font.bold = True; h2.text_frame.paragraphs[0].font.size = Pt(12); h2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; h2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    
    # Col 2 Text
    tf2 = col2.text_frame; tf2.word_wrap = True
    tf2.add_paragraph().text = "\n\n🔹 Problem: Users abandon carts because they fear the tedious returns process if the size is wrong.\n\n🔹 Solution: A policy allowing instant doorstep exchange for a different size, completely frictionless.\n\nKey Features:\n✅ One-Click Exchange – No approval needed for size swaps.\n✅ Instant Courier Dispatch – New size arrives same day.\n✅ Free Home Try-On – Try two sizes, return one.\n\n⭐ Impact: Reduces friction at checkout, but drastically increases reverse logistics costs for the business."
    for p in tf2.paragraphs: p.font.size = Pt(11); p.font.color.rgb = ajio_dark

    # Column 3: AJIO Fit Forums
    col3 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(0.8), Inches(3.1), Inches(6.3))
    col3.fill.solid(); col3.fill.fore_color.rgb = RGBColor(254, 240, 230); col3.line.color.rgb = ajio_dark; col3.line.dash_style = 4
    
    # Col 3 Header
    h3 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(0.9), Inches(2.7), Inches(0.4))
    h3.fill.solid(); h3.fill.fore_color.rgb = RGBColor(239, 68, 68); h3.line.color.rgb = RGBColor(239, 68, 68)
    h3.text_frame.paragraphs[0].text = "AJIO Fit Forums"; h3.text_frame.paragraphs[0].font.bold = True; h3.text_frame.paragraphs[0].font.size = Pt(12); h3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; h3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    
    # Col 3 Text
    tf3 = col3.text_frame; tf3.word_wrap = True
    tf3.add_paragraph().text = "\n\n🔹 Problem: Users leave the app to hunt for Reddit/YouTube reviews to see how garments fit on real people.\n\n🔹 Solution: An in-app social feed where users upload photos of themselves wearing the garment with their body stats.\n\nKey Features:\n✅ Shoppable UGC – See real users wearing the item.\n✅ Body Filter – Filter photos by your exact height/weight.\n✅ Creator Rewards – Points for uploading fit pics.\n\n⭐ Impact: High social proof, but suffers from the 'cold start' problem (needs massive user participation)."
    for p in tf3.paragraphs: p.font.size = Pt(11); p.font.color.rgb = ajio_dark

    # --- Slide 6: Prioritisation & Growth Loop (Based on Screenshot) ---
    slide6 = prs.slides.add_slide(blank_layout)
    
    # Title
    t_bg6 = slide6.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(4), Inches(0.5))
    t_bg6.text_frame.paragraphs[0].text = "Prioritisation & Growth Loop"; t_bg6.text_frame.paragraphs[0].font.bold = True; t_bg6.text_frame.paragraphs[0].font.size = Pt(20); t_bg6.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Red Line extending from title
    slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(0.4), Inches(5.2), Inches(0.03)).fill.solid()

    # --- LEFT COLUMN (Width: 5.8) ---
    # Prioritisation Header
    ph = slide6.shapes.add_textbox(Inches(0.2), Inches(0.8), Inches(3), Inches(0.4))
    ph.text_frame.paragraphs[0].text = "Prioritisation"; ph.text_frame.paragraphs[0].font.bold = True; ph.text_frame.paragraphs[0].font.size = Pt(16); ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    # Prioritisation Table
    rows, cols = 4, 7
    left, top, width, height = Inches(0.2), Inches(1.3), Inches(5.8), Inches(1.8)
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.4)
    for i in range(1, 7): table.columns[i].width = Inches(0.73)
    
    headers = ["Initiative", "Awareness", "Pain Point\nResolution", "Engagement\nImpact", "Effort", "Risk", "Priority\nScore"]
    for i in range(7):
        cell = table.cell(0, i)
        cell.text = headers[i]
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255,255,255)
        
    data = [
        ["1. AJIO Fit-Match", "4", "5", "4", "2", "2", "17 ✅"],
        ["2. Size-Swap", "4", "5", "2", "4", "5", "14"],
        ["3. Fit Forums", "3", "3", "4", "4", "3", "13"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if row_idx == 0: # Highlight winning row
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 245, 204)
                if col_idx == 0: cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68); cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                if col_idx == 0: cell.text_frame.paragraphs[0].font.bold = True

    # Growth Loop Header
    gh = slide6.shapes.add_textbox(Inches(0.2), Inches(3.3), Inches(3), Inches(0.4))
    gh.text_frame.paragraphs[0].text = "Growth Loop"; gh.text_frame.paragraphs[0].font.bold = True; gh.text_frame.paragraphs[0].font.size = Pt(16); gh.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    # Growth Loop Text
    gl = slide6.shapes.add_textbox(Inches(0.2), Inches(3.8), Inches(5.8), Inches(3.5))
    tf_gl = gl.text_frame; tf_gl.word_wrap = True
    
    p_g1 = tf_gl.paragraphs[0]; p_g1.text = "1 Discovery: Promote Fit-Match widget via pulse animations on high-abandonment SKU pages."; p_g1.font.bold = True; p_g1.font.size = Pt(11)
    p_g2 = tf_gl.add_paragraph(); p_g2.text = "2 Adoption: Guided onboarding prompts users to input their height and weight. First-time users see immediate visual try-on results."; p_g2.font.bold = True; p_g2.font.size = Pt(11)
    p_g3 = tf_gl.add_paragraph(); p_g3.text = "3 Positive Experience: User receives a highly-confident size recommendation, removing checkout anxiety and closing the trust gap."; p_g3.font.bold = True; p_g3.font.size = Pt(11)
    p_g4 = tf_gl.add_paragraph(); p_g4.text = "4 Data Flywheel: User purchases, keeps the item (no return), and leaves a verified review, which feeds back into the RAG model to help future users."; p_g4.font.bold = True; p_g4.font.size = Pt(11)
    p_g5 = tf_gl.add_paragraph(); p_g5.text = "5 Retention: User trusts AJIO's sizing unconditionally, increasing lifetime value and preventing platform switching."; p_g5.font.bold = True; p_g5.font.size = Pt(11)

    # --- RIGHT COLUMN (Width: 3.5, Peach Box) ---
    feat_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(0.8), Inches(3.5), Inches(6.3))
    feat_box.fill.solid(); feat_box.fill.fore_color.rgb = RGBColor(254, 240, 230); feat_box.line.color.rgb = RGBColor(254, 240, 230)
    
    tf_fb = feat_box.text_frame; tf_fb.word_wrap = True
    p_fb1 = tf_fb.paragraphs[0]; p_fb1.text = "Key Features (AJIO Fit-Match)\n"; p_fb1.font.bold = True; p_fb1.font.size = Pt(16); p_fb1.font.color.rgb = RGBColor(239, 68, 68)
    
    p_fb2 = tf_fb.add_paragraph(); p_fb2.text = "🚀 Fit-Match AI Engine – Get spot-on size recommendations based on your height, weight, and 1,000+ synthesized reviews. No more guessing!\n"
    p_fb2.font.bold = True; p_fb2.font.size = Pt(11); p_fb2.font.color.rgb = ajio_dark
    
    p_fb3 = tf_fb.add_paragraph(); p_fb3.text = "📸 Virtual Try-On – Visually map the garment onto a digital avatar to see how it hangs on your unique body type.\n"
    p_fb3.font.bold = True; p_fb3.font.size = Pt(11); p_fb3.font.color.rgb = ajio_dark
    
    p_fb4 = tf_fb.add_paragraph(); p_fb4.text = "🔔 Confidence Nudges – '80% of users your size bought a Medium and loved it!' Instant social proof at checkout.\n"
    p_fb4.font.bold = True; p_fb4.font.size = Pt(11); p_fb4.font.color.rgb = ajio_dark
    
    p_fb5 = tf_fb.add_paragraph(); p_fb5.text = "💬 Community Q&A (RAG) – Just ask: 'Is the fabric stretchy?' and Fit-Match will instantly summarize the answers from past buyers.\n"
    p_fb5.font.bold = True; p_fb5.font.size = Pt(11); p_fb5.font.color.rgb = ajio_dark
    
    p_fb6 = tf_fb.add_paragraph(); p_fb6.text = "📱 Seamless In-App Widget – Never leave the AJIO app to hunt for Reddit reviews again. Everything is right on the product page."
    p_fb6.font.bold = True; p_fb6.font.size = Pt(11); p_fb6.font.color.rgb = ajio_dark

    # --- Slide 7: Wireframing (AJIO Specific UI) ---
    slide7 = prs.slides.add_slide(blank_layout)
    
    # Title
    t_bg7 = slide7.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(3), Inches(0.5))
    t_bg7.text_frame.paragraphs[0].text = "Wireframing"; t_bg7.text_frame.paragraphs[0].font.bold = True; t_bg7.text_frame.paragraphs[0].font.size = Pt(20); t_bg7.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Red Line extending from title
    slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(0.4), Inches(7.5), Inches(0.03)).fill.solid()

    # --- LEFT COLUMN: AJIO Product Page (PDP) ---
    h_l = slide7.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(2), Inches(0.3))
    h_l.text_frame.paragraphs[0].text = "AJIO Product Page"; h_l.text_frame.paragraphs[0].font.bold = True; h_l.text_frame.paragraphs[0].font.size = Pt(14)
    desc_l = slide7.shapes.add_textbox(Inches(0.2), Inches(1.1), Inches(2.5), Inches(0.6))
    desc_l.text_frame.paragraphs[0].text = "The AI Fit-Match entry point intercepts the user right at the size selection friction point."; desc_l.text_frame.paragraphs[0].font.size = Pt(10); desc_l.text_frame.word_wrap = True

    # Mobile Screen 1 (AJIO UI)
    m1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(1.8), Inches(3.8))
    m1.fill.solid(); m1.fill.fore_color.rgb = RGBColor(255, 255, 255); m1.line.color.rgb = RGBColor(200, 200, 200)
    
    # AJIO Top Nav Bar
    nav1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(1.8), Inches(0.25))
    nav1.fill.solid(); nav1.fill.fore_color.rgb = RGBColor(255, 255, 255); nav1.line.color.rgb = RGBColor(230, 230, 230)
    nav_txt = nav1.text_frame.paragraphs[0]; nav_txt.text = "≡    AJIO    🔍 🛍"; nav_txt.font.bold = True; nav_txt.font.size = Pt(9); nav_txt.font.color.rgb = ajio_dark
    
    # Product Image Area
    slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.05), Inches(1.8), Inches(1.8)).fill.solid()
    
    # Product Details (Brand, Title, Price)
    b_txt = slide7.shapes.add_textbox(Inches(0.55), Inches(3.9), Inches(1.0), Inches(0.2))
    b_txt.text_frame.paragraphs[0].text = "GAP"; b_txt.text_frame.paragraphs[0].font.bold = True; b_txt.text_frame.paragraphs[0].font.size = Pt(10)
    t_txt = slide7.shapes.add_textbox(Inches(0.55), Inches(4.05), Inches(1.6), Inches(0.2))
    t_txt.text_frame.paragraphs[0].text = "Men's Graphic Print Hoodie"; t_txt.text_frame.paragraphs[0].font.size = Pt(8)
    p_txt = slide7.shapes.add_textbox(Inches(0.55), Inches(4.2), Inches(1.0), Inches(0.2))
    p_txt.text_frame.paragraphs[0].text = "₹1,499"; p_txt.text_frame.paragraphs[0].font.bold = True; p_txt.text_frame.paragraphs[0].font.size = Pt(10)
    
    # Fit-Match Button UI (Highlighted right above sizes)
    btn1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.5), Inches(1.6), Inches(0.25))
    btn1.fill.solid(); btn1.fill.fore_color.rgb = RGBColor(218, 227, 243); btn1.line.color.rgb = ajio_dark
    b1_txt = btn1.text_frame.paragraphs[0]; b1_txt.text = "✨ Try AI Fit-Match"; b1_txt.font.bold = True; b1_txt.font.size = Pt(8); b1_txt.font.color.rgb = ajio_dark
    
    # Dashed Pink Highlight Box
    dh1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(4.45), Inches(1.7), Inches(0.35))
    dh1.fill.background(); dh1.line.color.rgb = RGBColor(255, 20, 147); dh1.line.dash_style = 4; dh1.line.width = Pt(1.5)
    
    # Pointers
    ptr1 = slide7.shapes.add_textbox(Inches(2.4), Inches(4.35), Inches(1.8), Inches(0.6))
    ptr1.text_frame.paragraphs[0].text = "⬅ Replaces standard size chart with confident AI."; ptr1.text_frame.paragraphs[0].font.size = Pt(9); ptr1.text_frame.word_wrap = True

    # Standard Size Selector UI
    s_lbl = slide7.shapes.add_textbox(Inches(0.55), Inches(4.75), Inches(1.0), Inches(0.2))
    s_lbl.text_frame.paragraphs[0].text = "Select Size"; s_lbl.text_frame.paragraphs[0].font.size = Pt(8)
    slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(4.95), Inches(0.2), Inches(0.2)).fill.solid() # S
    slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(4.95), Inches(0.2), Inches(0.2)).fill.solid() # M
    slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(4.95), Inches(0.2), Inches(0.2)).fill.solid() # L
    
    # Add to Bag / Wishlist Footer
    add_btn = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.3), Inches(1.8), Inches(0.3))
    add_btn.fill.solid(); add_btn.fill.fore_color.rgb = ajio_dark
    a_txt = add_btn.text_frame.paragraphs[0]; a_txt.text = "ADD TO BAG"; a_txt.font.size = Pt(8); a_txt.font.bold = True; a_txt.font.color.rgb = RGBColor(255,255,255)

    # --- RIGHT AREA: AJIO AI Fit-Match Flow ---
    h_r = slide7.shapes.add_textbox(Inches(4.5), Inches(0.8), Inches(4), Inches(0.3))
    h_r.text_frame.paragraphs[0].text = "AI Sizing Assistant: In-App Flow"; h_r.text_frame.paragraphs[0].font.bold = True; h_r.text_frame.paragraphs[0].font.size = Pt(14)

    # Mobile Screen 2 (Input Modal)
    m2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1.3), Inches(1.7), Inches(3.5))
    m2.fill.solid(); m2.fill.fore_color.rgb = RGBColor(255, 255, 255); m2.line.color.rgb = RGBColor(200, 200, 200)
    
    # Nav
    nav2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(1.3), Inches(1.7), Inches(0.25))
    nav2.fill.solid(); nav2.fill.fore_color.rgb = RGBColor(255, 255, 255); nav2.line.color.rgb = RGBColor(230, 230, 230)
    nav_txt2 = nav2.text_frame.paragraphs[0]; nav_txt2.text = "←  AJIO Fit-Match"; nav_txt2.font.bold = True; nav_txt2.font.size = Pt(9); nav_txt2.font.color.rgb = ajio_dark
    
    # Inputs
    inp1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(2.0), Inches(1.5), Inches(0.25))
    inp1.fill.solid(); inp1.fill.fore_color.rgb = RGBColor(250, 250, 250); inp1.line.color.rgb = RGBColor(200,200,200)
    inp1.text_frame.paragraphs[0].text = "Height (cm)"; inp1.text_frame.paragraphs[0].font.size = Pt(8); inp1.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,100,100)
    
    inp2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(2.4), Inches(1.5), Inches(0.25))
    inp2.fill.solid(); inp2.fill.fore_color.rgb = RGBColor(250, 250, 250); inp2.line.color.rgb = RGBColor(200,200,200)
    inp2.text_frame.paragraphs[0].text = "Weight (kg)"; inp2.text_frame.paragraphs[0].font.size = Pt(8); inp2.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,100,100)
    
    btn2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(3.0), Inches(1.1), Inches(0.3))
    btn2.fill.solid(); btn2.fill.fore_color.rgb = ajio_dark
    btn2.text_frame.paragraphs[0].text = "Find Size"; btn2.text_frame.paragraphs[0].font.size = Pt(8); btn2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    
    dh2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.85), Inches(1.9), Inches(1.6), Inches(1.5))
    dh2.fill.background(); dh2.line.color.rgb = RGBColor(255, 20, 147); dh2.line.dash_style = 4; dh2.line.width = Pt(1.5)

    # Mobile Screen 3 (RAG Synthesis / Recommendation)
    m3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.3), Inches(1.7), Inches(3.5))
    m3.fill.solid(); m3.fill.fore_color.rgb = RGBColor(255, 255, 255); m3.line.color.rgb = RGBColor(200, 200, 200)
    
    nav3 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(1.3), Inches(1.7), Inches(0.25))
    nav3.fill.solid(); nav3.fill.fore_color.rgb = RGBColor(255, 255, 255); nav3.line.color.rgb = RGBColor(230, 230, 230)
    nav_txt3 = nav3.text_frame.paragraphs[0]; nav_txt3.text = "←  Your Fit"; nav_txt3.font.bold = True; nav_txt3.font.size = Pt(9); nav_txt3.font.color.rgb = ajio_dark
    
    rec_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.7), Inches(1.5), Inches(1.0))
    rec_box.fill.solid(); rec_box.fill.fore_color.rgb = RGBColor(220, 252, 231); rec_box.line.color.rgb = RGBColor(34, 197, 94)
    r_txt = rec_box.text_frame
    r_txt.paragraphs[0].text = "Recommended: L"; r_txt.paragraphs[0].font.bold = True; r_txt.paragraphs[0].font.size = Pt(10); r_txt.paragraphs[0].font.color.rgb = RGBColor(22, 101, 52)
    r_txt.add_paragraph().text = "\n98% Match based on 1,240 verified reviews."; r_txt.paragraphs[1].font.size = Pt(7); r_txt.paragraphs[1].font.color.rgb = ajio_dark
    
    trust_box = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(2.9), Inches(1.5), Inches(0.5))
    trust_box.fill.solid(); trust_box.fill.fore_color.rgb = RGBColor(250, 250, 250); trust_box.line.color.rgb = RGBColor(200,200,200)
    t_txt = trust_box.text_frame.paragraphs[0]; t_txt.text = "👍 82% of buyers with your stats kept this size."; t_txt.font.size = Pt(7); t_txt.font.color.rgb = ajio_dark; t_txt.alignment = PP_ALIGN.CENTER
    
    dh3 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.85), Inches(1.6), Inches(1.6), Inches(1.2))
    dh3.fill.background(); dh3.line.color.rgb = RGBColor(255, 20, 147); dh3.line.dash_style = 4; dh3.line.width = Pt(1.5)

    # Mobile Screen 4 (Virtual Try-on)
    m4 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.3), Inches(1.7), Inches(3.5))
    m4.fill.solid(); m4.fill.fore_color.rgb = RGBColor(255, 255, 255); m4.line.color.rgb = RGBColor(200, 200, 200)
    
    nav4 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(1.3), Inches(1.7), Inches(0.25))
    nav4.fill.solid(); nav4.fill.fore_color.rgb = RGBColor(255, 255, 255); nav4.line.color.rgb = RGBColor(230, 230, 230)
    nav_txt4 = nav4.text_frame.paragraphs[0]; nav_txt4.text = "←  Visual Try-On"; nav_txt4.font.bold = True; nav_txt4.font.size = Pt(9); nav_txt4.font.color.rgb = ajio_dark
    
    vto_img = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.9), Inches(1.7), Inches(1.5), Inches(2.2))
    vto_img.fill.solid(); vto_img.fill.fore_color.rgb = RGBColor(240, 240, 240); vto_img.line.color.rgb = RGBColor(200, 200, 200)
    vto_img.text_frame.paragraphs[0].text = "👤\nDigital\nMannequin\nOverlay"; vto_img.text_frame.paragraphs[0].font.size = Pt(10); vto_img.text_frame.paragraphs[0].font.color.rgb = RGBColor(150,150,150); vto_img.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    btn4 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.9), Inches(4.1), Inches(1.5), Inches(0.3))
    btn4.fill.solid(); btn4.fill.fore_color.rgb = ajio_dark
    btn4_txt = btn4.text_frame.paragraphs[0]; btn4_txt.text = "PROCEED TO BAG"; btn4_txt.font.size = Pt(8); btn4_txt.font.bold = True; btn4_txt.font.color.rgb = RGBColor(255,255,255)
    
    dh4 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.85), Inches(1.6), Inches(1.6), Inches(2.4))
    dh4.fill.background(); dh4.line.color.rgb = RGBColor(255, 20, 147); dh4.line.dash_style = 4; dh4.line.width = Pt(1.5)

    # Descriptions under the flow
    d2 = slide7.shapes.add_textbox(Inches(3.8), Inches(4.9), Inches(1.8), Inches(1))
    d2.text_frame.paragraphs[0].text = "📐 Step 1: Input stats\nUser enters height & weight via in-app modal."; d2.text_frame.paragraphs[0].font.size = Pt(9); d2.text_frame.word_wrap = True
    
    d3 = slide7.shapes.add_textbox(Inches(5.8), Inches(4.9), Inches(1.8), Inches(1))
    d3.text_frame.paragraphs[0].text = "🤖 Step 2: AI Synthesis\nRAG Engine parses 1,000+ reviews for a high-confidence recommendation."; d3.text_frame.paragraphs[0].font.size = Pt(9); d3.text_frame.word_wrap = True
    
    d4 = slide7.shapes.add_textbox(Inches(7.8), Inches(4.9), Inches(1.8), Inches(1))
    d4.text_frame.paragraphs[0].text = "📸 Step 3: Try-On\nUser sees exact fit visualization, eliminating purchase hesitation."; d4.text_frame.paragraphs[0].font.size = Pt(9); d4.text_frame.word_wrap = True

    # --- Slide 8: Success Metrics ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Success Metrics: Measuring Impact")
    
    rows, cols = 5, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3)
    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Type", "Metric", "Calculation", "Rationale"]
    for i in range(4):
        cell = table.cell(0, i)
        cell.text = headers[i]
        cell.fill.solid(); cell.fill.fore_color.rgb = ajio_dark
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        
    data = [
        ["North Star", "Wishlist-to-Purchase Rate", "Purchases / Wishlist Adds", "Direct measurement of core business goal."],
        ["Leading", "Fit-Match Engagement Rate", "AI Widget Clicks / Views", "Proves the widget successfully intercepts doubt."],
        ["Leading", "Time-to-Purchase", "Time from Wishlist to Checkout", "Decrease indicates elimination of review-hunting."],
        ["Guardrail", "Size-Related Return Rate", "Size Returns / Total purchases", "Ensures AI recommendations are accurate. Must stay flat/decrease."]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if row_idx == 3: 
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 230, 230)

    # --- Slide 9: Testing & Rollout Plan ---
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Testing & Rollout Plan: A/B Testing Cohorts")
    
    tb = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
    tb.fill.solid(); tb.fill.fore_color.rgb = bg_light; tb.line.color.rgb = ajio_dark
    tb.text_frame.text = "Phased Rollout Strategy (A/B Test):\n\nPhase 1 (Pilot - Top 10% High-Volume SKUs):\n- Control Group (50%): Standard AJIO product page.\n- Variant Group (50%): Product page with 'Fit-Match AI' widget.\n\nEvaluation Criteria:\n- Monitor Time-to-Purchase and Widget Engagement over 14 days.\n- Wait 30 days to measure the Guardrail Metric (Size-Related Returns) to ensure the AI isn't causing poor purchases.\n\nPhase 2 (General Availability):\n- If Wishlist-to-Purchase increases by >3% and Returns remain flat, scale the RAG pipeline to all SKUs with >50 reviews."
    for p in tb.text_frame.paragraphs: p.font.size = Pt(14); p.font.color.rgb = ajio_dark

    # --- Slide 10: Risks & Mitigations ---
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Risks & Mitigations")
    
    r1 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3), Inches(3.5))
    r1.fill.solid(); r1.fill.fore_color.rgb = bg_light; r1.line.color.rgb = ajio_dark
    r1.text_frame.text = "Risk 1: LLM Hallucinations\n\nMitigation: Strict RAG boundaries. The AI only triggers if >5 highly relevant reviews exist for the SKU. Otherwise, it gracefully falls back."
    for p in r1.text_frame.paragraphs: p.font.size = Pt(14); p.font.color.rgb = ajio_dark

    r2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.75), Inches(1.5), Inches(2.5), Inches(3.5))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg_light; r2.line.color.rgb = ajio_dark
    r2.text_frame.text = "Risk 2: Low Discoverability\n\nMitigation: Implement micro-animations (glow effects) on the widget if the user lingers on the size-selector tool for >3 seconds."
    for p in r2.text_frame.paragraphs: p.font.size = Pt(14); p.font.color.rgb = ajio_dark

    # --- Slide 8: Launch Strategy and System Design (Based on Screenshot) ---
    slide8 = prs.slides.add_slide(blank_layout)
    
    # Title
    t_bg8 = slide8.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(5), Inches(0.5))
    t_bg8.text_frame.paragraphs[0].text = "Launch Strategy and System Design"; t_bg8.text_frame.paragraphs[0].font.bold = True; t_bg8.text_frame.paragraphs[0].font.size = Pt(20); t_bg8.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Red Line extending from title
    slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0.4), Inches(5.2), Inches(0.03)).fill.solid()

    # --- TOP LEFT: Launch Strategy ---
    ls_h = slide8.shapes.add_textbox(Inches(0.2), Inches(0.7), Inches(3), Inches(0.3))
    ls_h.text_frame.paragraphs[0].text = "Launch Strategy:"; ls_h.text_frame.paragraphs[0].font.bold = True; ls_h.text_frame.paragraphs[0].font.size = Pt(14); ls_h.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    ls_text = slide8.shapes.add_textbox(Inches(0.2), Inches(1.1), Inches(4.2), Inches(4.0))
    tf_ls = ls_text.text_frame; tf_ls.word_wrap = True
    
    p_ls1 = tf_ls.paragraphs[0]; p_ls1.text = "Target Audience"; p_ls1.font.bold = True; p_ls1.font.size = Pt(11)
    p_ls2 = tf_ls.add_paragraph(); p_ls2.text = "• Primary: Hesitant wishlisters with high cart abandonment.\n• Secondary: High-frequency shoppers looking for new brands."; p_ls2.font.size = Pt(10)
    
    p_ls3 = tf_ls.add_paragraph(); p_ls3.text = "\nValue Proposition"; p_ls3.font.bold = True; p_ls3.font.size = Pt(11)
    p_ls4 = tf_ls.add_paragraph(); p_ls4.text = "• AI-Powered Personalization: Exact size recommendations via RAG.\n• Virtual Try-On: Visual confidence mapping garment to body type.\n• Higher Conversion: Proactive sizing nudges to reduce drop-off."; p_ls4.font.size = Pt(10)
    
    p_ls5 = tf_ls.add_paragraph(); p_ls5.text = "\nPromotion Channels"; p_ls5.font.bold = True; p_ls5.font.size = Pt(11)
    p_ls6 = tf_ls.add_paragraph(); p_ls6.text = "• In-App Notifications: 'Find your perfect fit instantly.'\n• Wishlist Nudges: Targeted pushes for abandoned items.\n• Checkout Real-Estate: Pulse animations near the size selector."; p_ls6.font.size = Pt(10)

    # --- TOP RIGHT: High Level Design Diagram ---
    hld_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(0.8), Inches(5.0), Inches(4.2))
    hld_box.fill.solid(); hld_box.fill.fore_color.rgb = RGBColor(224, 242, 254); hld_box.line.color.rgb = RGBColor(224, 242, 254)
    
    hld_title = slide8.shapes.add_textbox(Inches(4.7), Inches(0.9), Inches(2), Inches(0.3))
    hld_title.text_frame.paragraphs[0].text = "High Level Design"; hld_title.text_frame.paragraphs[0].font.bold = True; hld_title.text_frame.paragraphs[0].font.size = Pt(12); hld_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    # Drawing the Architecture Blocks
    # User
    u_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(2.3), Inches(0.6), Inches(0.4))
    u_box.fill.solid(); u_box.fill.fore_color.rgb = RGBColor(200,200,200); u_box.text_frame.paragraphs[0].text = "👤 User"; u_box.text_frame.paragraphs[0].font.size = Pt(9); u_box.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # AJIO App
    app_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(2.3), Inches(1.0), Inches(0.4))
    app_box.fill.solid(); app_box.fill.fore_color.rgb = RGBColor(255,255,255); app_box.line.color.rgb = ajio_dark; app_box.text_frame.paragraphs[0].text = "📱 AJIO App"; app_box.text_frame.paragraphs[0].font.size = Pt(9); app_box.text_frame.paragraphs[0].font.color.rgb = ajio_dark

    # AI Engine (RAG)
    ai_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.4), Inches(1.0), Inches(0.5))
    ai_box.fill.solid(); ai_box.fill.fore_color.rgb = RGBColor(255,228,225); ai_box.line.color.rgb = RGBColor(239,68,68); ai_box.text_frame.paragraphs[0].text = "🧠 AI Engine\n(RAG)"; ai_box.text_frame.paragraphs[0].font.size = Pt(8); ai_box.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Context Handler
    ctx_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(2.2), Inches(1.0), Inches(0.5))
    ctx_box.fill.solid(); ctx_box.fill.fore_color.rgb = RGBColor(255,228,225); ctx_box.line.color.rgb = RGBColor(239,68,68); ctx_box.text_frame.paragraphs[0].text = "⚙️ Context\nHandler"; ctx_box.text_frame.paragraphs[0].font.size = Pt(8); ctx_box.text_frame.paragraphs[0].font.color.rgb = ajio_dark

    # Virtual Try-On Service (VTO)
    vto_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(3.0), Inches(1.0), Inches(0.5))
    vto_box.fill.solid(); vto_box.fill.fore_color.rgb = RGBColor(220,252,231); vto_box.line.color.rgb = RGBColor(34,197,94); vto_box.text_frame.paragraphs[0].text = "📸 VTO\nService"; vto_box.text_frame.paragraphs[0].font.size = Pt(8); vto_box.text_frame.paragraphs[0].font.color.rgb = ajio_dark

    # DB/Storage
    db_box = slide8.shapes.add_shape(MSO_SHAPE.CAN, Inches(8.6), Inches(1.3), Inches(0.8), Inches(2.3))
    db_box.fill.solid(); db_box.fill.fore_color.rgb = RGBColor(59, 130, 246); db_box.line.color.rgb = RGBColor(29, 78, 216)
    d_txt = slide8.shapes.add_textbox(Inches(8.5), Inches(0.9), Inches(1.0), Inches(0.3))
    d_txt.text_frame.paragraphs[0].text = "Vector DB"; d_txt.text_frame.paragraphs[0].font.size = Pt(8); d_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Lines
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(2.5), Inches(5.8), Inches(2.5)) # User to App
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.5), Inches(7.2), Inches(1.6)) # App to AI
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.5), Inches(7.2), Inches(2.4)) # App to Context
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.5), Inches(7.2), Inches(3.2)) # App to VTO
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(1.6), Inches(8.6), Inches(1.6)) # AI to DB
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(2.4), Inches(8.6), Inches(2.4)) # Context to DB
    slide8.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(3.2), Inches(8.6), Inches(3.2)) # VTO to DB

    # DB Bullet Points
    db_pts = slide8.shapes.add_textbox(Inches(7.2), Inches(3.8), Inches(2.5), Inches(1.0))
    db_pts.text_frame.paragraphs[0].text = "• AI Training Data (Reviews & Images)\n• Vector Cache (Redis)\n• User Context DB"; db_pts.text_frame.paragraphs[0].font.size = Pt(8)

    # --- BOTTOM ROW: Launch Timeline ---
    lt_h = slide8.shapes.add_textbox(Inches(0.2), Inches(5.2), Inches(3), Inches(0.3))
    lt_h.text_frame.paragraphs[0].text = "Launch Timeline"; lt_h.text_frame.paragraphs[0].font.bold = True; lt_h.text_frame.paragraphs[0].font.size = Pt(14); lt_h.text_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    # 4 Columns for phases
    w = Inches(2.3)
    y_pos = Inches(5.6)
    
    # Phase 1
    p1 = slide8.shapes.add_textbox(Inches(0.2), y_pos, w, Inches(1.5)); tf_p1 = p1.text_frame; tf_p1.word_wrap = True
    tf_p1.paragraphs[0].text = "Phase 1: Pre-Launch\n(1 Week Before)"; tf_p1.paragraphs[0].font.bold = True; tf_p1.paragraphs[0].font.size = Pt(11); tf_p1.paragraphs[0].font.color.rgb = RGBColor(128,0,128)
    tf_p1.add_paragraph().text = "✅ Beta test Fit-Match AI with select power users.\n✅ Tease via emails and in-app notifications."; tf_p1.paragraphs[1].font.size = Pt(10); tf_p1.paragraphs[1].font.bold = True
    
    # Phase 2
    p2 = slide8.shapes.add_textbox(Inches(2.6), y_pos, w, Inches(1.5)); tf_p2 = p2.text_frame; tf_p2.word_wrap = True
    tf_p2.paragraphs[0].text = "Phase 2: Launch Week*"; tf_p2.paragraphs[0].font.bold = True; tf_p2.paragraphs[0].font.size = Pt(11); tf_p2.paragraphs[0].font.color.rgb = RGBColor(128,0,128)
    tf_p2.add_paragraph().text = "✅ Limited release to Top 10% high-abandonment SKUs.\n✅ A/B test UI widget placement."; tf_p2.paragraphs[1].font.size = Pt(10); tf_p2.paragraphs[1].font.bold = True
    
    # Phase 3
    p3 = slide8.shapes.add_textbox(Inches(5.0), y_pos, w, Inches(1.5)); tf_p3 = p3.text_frame; tf_p3.word_wrap = True
    tf_p3.paragraphs[0].text = "Phase 3: Post-Launch\n(2-4 Weeks)"; tf_p3.paragraphs[0].font.bold = True; tf_p3.paragraphs[0].font.size = Pt(11); tf_p3.paragraphs[0].font.color.rgb = RGBColor(128,0,128)
    tf_p3.add_paragraph().text = "✅ Global launch across all clothing categories.\n✅ Incentivize users to leave detailed sizing reviews."; tf_p3.paragraphs[1].font.size = Pt(10); tf_p3.paragraphs[1].font.bold = True
    
    # Phase 4
    p4 = slide8.shapes.add_textbox(Inches(7.4), y_pos, w, Inches(1.5)); tf_p4 = p4.text_frame; tf_p4.word_wrap = True
    tf_p4.paragraphs[0].text = "Phase 4: Post-Launch\n(Month 2+)"; tf_p4.paragraphs[0].font.bold = True; tf_p4.paragraphs[0].font.size = Pt(11); tf_p4.paragraphs[0].font.color.rgb = RGBColor(128,0,128)
    tf_p4.add_paragraph().text = "✅ AI-driven engagement nudges at checkout.\n✅ Continuous RAG model training via user return data."; tf_p4.paragraphs[1].font.size = Pt(10); tf_p4.paragraphs[1].font.bold = True
    # --- Slide 9: Metrics, Pitfall & Mitigation (Based on Screenshot) ---
    slide9 = prs.slides.add_slide(blank_layout)
    
    # Title
    t_bg9 = slide9.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(4), Inches(0.5))
    t_bg9.text_frame.paragraphs[0].text = "Metrics, Pitfall & Mitigation"; t_bg9.text_frame.paragraphs[0].font.bold = True; t_bg9.text_frame.paragraphs[0].font.size = Pt(20); t_bg9.text_frame.paragraphs[0].font.color.rgb = ajio_dark
    
    # Red Line extending from title
    slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(0.4), Inches(5.4), Inches(0.03)).fill.solid()

    # --- LEFT COLUMN: Metrics Table ---
    mh = slide9.shapes.add_textbox(Inches(0.2), Inches(0.8), Inches(4), Inches(0.3))
    mh.text_frame.paragraphs[0].text = "How to measure success?"; mh.text_frame.paragraphs[0].font.bold = True; mh.text_frame.paragraphs[0].font.size = Pt(16)

    rows, cols = 6, 4
    left, top, width, height = Inches(0.2), Inches(1.3), Inches(5.8), Inches(3.8)
    table_shape = slide9.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(2.0)
    
    headers = ["METRIC", "ASPECT", "TYPE", "OBJECTIVE"]
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = headers[i]
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(244, 164, 196) # Pink header
        
    data = [
        ["Wishlist-to-Purchase Conversion Rate", "North-Star\nMetric (NSM)", "Conversion", "Measures overall feature impact on closing the visual trust gap."],
        ["Size-Related Return Rate", "Primary /\nGuardrail", "Cost\nSavings", "Track reduction of returns specifically due to 'wrong size'."],
        ["Fit-Match Widget CTR", "Leading - L1", "Engagement", "Measures adoption and click-through of the AI sizing tool."],
        ["Virtual Try-On Completion", "Leading - L1", "Feature\nUsage", "Measures how many users complete the visual mapping process."],
        ["Time-to-Checkout", "Supporting\nMetrics (L2)", "Friction", "Measure reduction in time spent hesitating on the product page."]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            if col_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(214, 134, 166) # Darker pink
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # Add thin border manually using lines or just rely on default
            
    # --- RIGHT COLUMN: Pitfall & Mitigation ---
    ph = slide9.shapes.add_textbox(Inches(6.2), Inches(0.8), Inches(3.5), Inches(0.3))
    ph.text_frame.paragraphs[0].text = "Pitfall & Mitigation"; ph.text_frame.paragraphs[0].font.bold = True; ph.text_frame.paragraphs[0].font.size = Pt(16)

    # Box 1
    pb1 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(1.3), Inches(3.5), Inches(1.0))
    pb1.fill.solid(); pb1.fill.fore_color.rgb = RGBColor(255,255,255); pb1.line.color.rgb = RGBColor(200,200,200)
    pt1 = pb1.text_frame; pt1.word_wrap = True
    pt1.paragraphs[0].text = "AI Hallucinations (Bad Sizing)"; pt1.paragraphs[0].font.bold = True; pt1.paragraphs[0].font.size = Pt(9); pt1.paragraphs[0].font.color.rgb = ajio_dark
    p1_1 = pt1.add_paragraph(); p1_1.text = "Pitfall: AI suggests wrong size, leading to returns."; p1_1.font.size = Pt(8); p1_1.font.bold = True; p1_1.font.color.rgb = ajio_dark
    p1_2 = pt1.add_paragraph(); p1_2.text = "Mitigation: Confidence Threshold. Fallback to standard chart if AI consensus is <80%."; p1_2.font.size = Pt(8); p1_2.font.color.rgb = ajio_dark

    # Box 2
    pb2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(2.4), Inches(3.5), Inches(1.0))
    pb2.fill.solid(); pb2.fill.fore_color.rgb = RGBColor(255,255,255); pb2.line.color.rgb = RGBColor(200,200,200)
    pt2 = pb2.text_frame; pt2.word_wrap = True
    pt2.paragraphs[0].text = "Cold Start Problem"; pt2.paragraphs[0].font.bold = True; pt2.paragraphs[0].font.size = Pt(9); pt2.paragraphs[0].font.color.rgb = ajio_dark
    p2_1 = pt2.add_paragraph(); p2_1.text = "Pitfall: New SKUs have no reviews to generate sizes."; p2_1.font.size = Pt(8); p2_1.font.bold = True; p2_1.font.color.rgb = ajio_dark
    p2_2 = pt2.add_paragraph(); p2_2.text = "Mitigation: AI aggregates reviews from similar silhouettes within the same brand for a baseline."; p2_2.font.size = Pt(8); p2_2.font.color.rgb = ajio_dark

    # Box 3
    pb3 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(3.5), Inches(3.5), Inches(1.0))
    pb3.fill.solid(); pb3.fill.fore_color.rgb = RGBColor(255,255,255); pb3.line.color.rgb = RGBColor(200,200,200)
    pt3 = pb3.text_frame; pt3.word_wrap = True
    pt3.paragraphs[0].text = "Low User Adoption"; pt3.paragraphs[0].font.bold = True; pt3.paragraphs[0].font.size = Pt(9); pt3.paragraphs[0].font.color.rgb = ajio_dark
    p3_1 = pt3.add_paragraph(); p3_1.text = "Pitfall: Users ignore the Fit-Match AI button."; p3_1.font.size = Pt(8); p3_1.font.bold = True; p3_1.font.color.rgb = ajio_dark
    p3_2 = pt3.add_paragraph(); p3_2.text = "Mitigation: Implement micro-animations (glow effects) near the size-selector to draw attention."; p3_2.font.size = Pt(8); p3_2.font.color.rgb = ajio_dark

    # Box 4
    pb4 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(4.6), Inches(3.5), Inches(1.0))
    pb4.fill.solid(); pb4.fill.fore_color.rgb = RGBColor(255,255,255); pb4.line.color.rgb = RGBColor(200,200,200)
    pt4 = pb4.text_frame; pt4.word_wrap = True
    pt4.paragraphs[0].text = "High Latency (Slow RAG)"; pt4.paragraphs[0].font.bold = True; pt4.paragraphs[0].font.size = Pt(9); pt4.paragraphs[0].font.color.rgb = ajio_dark
    p4_1 = pt4.add_paragraph(); p4_1.text = "Pitfall: Users abandon cart if AI takes too long to load."; p4_1.font.size = Pt(8); p4_1.font.bold = True; p4_1.font.color.rgb = ajio_dark
    p4_2 = pt4.add_paragraph(); p4_2.text = "Mitigation: Pre-compute and cache sizing summaries for the top 20% of best-selling SKUs."; p4_2.font.size = Pt(8); p4_2.font.color.rgb = ajio_dark

    prs.save('AJIO_Final_Deck_V20.pptx')
    print("Created AJIO_Final_Deck_V20.pptx successfully (Fixed Actors on Slide 1).")

if __name__ == '__main__':
    create_visual_deck()
